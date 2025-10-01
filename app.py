# app.py
import streamlit as st
import pandas as pd
import numpy as np
import os
import re
import pickle
from io import BytesIO
from typing import List, Dict

# file parsers
from PyPDF2 import PdfReader
import docx
from pptx import Presentation

# OpenAI
import openai

# try FAISS, else try Annoy, else fallback to sklearn brute force
FAISS_AVAILABLE = False
ANNOY_AVAILABLE = False
try:
    import faiss
    FAISS_AVAILABLE = True
except Exception:
    try:
        from annoy import AnnoyIndex
        ANNOY_AVAILABLE = True
    except Exception:
        from sklearn.metrics.pairwise import linear_kernel
        from sklearn.feature_extraction.text import TfidfVectorizer

# --- Streamlit config ---
st.set_page_config(page_title="Aircraft Faults — Embeddings + FAISS", layout="wide")
st.title("✈️ Tra cứu & Phân tích sự cố tàu bay — Embeddings + FAISS")

# -------------------------
# Sidebar: settings & keys
# -------------------------
with st.sidebar:
    st.header("Cấu hình / Settings")
    st.markdown("**OpenAI API key**: (tốt nhất set ở Secrets trên Streamlit Cloud thay vì nhập ở đây)")
    key_input = st.text_input("Hoặc dán API key (tạm) — bắt đầu bằng sk-", type="password")
    DEFAULT_EMBED_MODEL = st.selectbox("Embedding model", options=["text-embedding-3-small", "text-embedding-3-large"], index=0)
    DEFAULT_CHAT_MODEL = st.selectbox("Chat model", options=["gpt-4o-mini","gpt-4o","gpt-3.5-turbo"], index=0)
    chunk_size = st.number_input("Chunk size (chars)", value=1200, min_value=200)
    chunk_overlap = st.number_input("Chunk overlap (chars)", value=200, min_value=0)
    top_k = st.slider("Số chunk lấy làm context (top k)", 1, 20, 6)
    use_faiss = st.checkbox("Dùng FAISS nếu có (fallback: Annoy hoặc TF-IDF)", value=True)
    st.markdown("---")
    st.markdown("**Lưu ý chi phí & bảo mật**:\n- Embeddings + LLM gọi tới OpenAI tốn tiền. Giữ API key an toàn.\n- Không upload dữ liệu nhạy cảm lên cloud nếu không muốn chia sẻ.")

# Get API key: prefer st.secrets (Streamlit Cloud) else env var else text input
OPENAI_KEY = None
if st.secrets and "OPENAI_API_KEY" in st.secrets:
    OPENAI_KEY = st.secrets["OPENAI_API_KEY"]
elif os.environ.get("OPENAI_API_KEY"):
    OPENAI_KEY = os.environ.get("OPENAI_API_KEY")
elif key_input:
    OPENAI_KEY = key_input

if OPENAI_KEY:
    openai.api_key = OPENAI_KEY
else:
    st.warning("Chưa có OpenAI API key — app sẽ dùng chế độ fallback (TF-IDF) nếu available.")

# -------------------------
# Helpers: read files
# -------------------------
def read_pdf(file_bytes: BytesIO) -> str:
    text = ""
    try:
        reader = PdfReader(file_bytes)
        for p in reader.pages:
            page_text = p.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.warning(f"PDF parse lỗi: {e}")
    return text

def read_docx(file_bytes: BytesIO) -> str:
    text = ""
    try:
        doc = docx.Document(file_bytes)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        st.warning(f"DOCX parse lỗi: {e}")
    return text

def read_pptx(file_bytes: BytesIO) -> str:
    text = ""
    try:
        prs = Presentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        st.warning(f"PPTX parse lỗi: {e}")
    return text

def read_table_like(file_bytes: BytesIO, filename: str) -> str:
    text = ""
    try:
        if filename.lower().endswith(".csv"):
            df = pd.read_csv(file_bytes, dtype=str, keep_default_na=False)
            text = df.astype(str).to_csv(index=False)
        else:
            df = pd.read_excel(file_bytes, dtype=str)
            text = df.astype(str).to_csv(index=False)
    except Exception as e:
        st.warning(f"Excel/CSV parse lỗi: {e}")
    return text

def read_text_file(file_bytes: BytesIO) -> str:
    try:
        return file_bytes.read().decode("utf-8", errors="ignore")
    except Exception:
        return ""

def extract_text_from_file(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    bio = BytesIO(uploaded_file.read())
    bio.seek(0)
    if name.endswith(".pdf"):
        return read_pdf(bio)
    elif name.endswith(".docx"):
        bio.seek(0)
        return read_docx(bio)
    elif name.endswith(".pptx"):
        bio.seek(0)
        return read_pptx(bio)
    elif name.endswith(".csv") or name.endswith(".xls") or name.endswith(".xlsx"):
        bio.seek(0)
        return read_table_like(bio, name)
    else:
        # fallback: txt / unknown
        bio.seek(0)
        return read_text_file(bio)

# -------------------------
# Chunking
# -------------------------
def chunk_text(text: str, chunk_size:int=1200, overlap:int=200) -> List[str]:
    text = re.sub(r'\s+', ' ', text).strip()
    if not text:
        return []
    chunks = []
    start = 0
    L = len(text)
    while start < L:
        end = min(L, start + chunk_size)
        # try to break at nearest sentence end before end (up to 80 chars)
        if end < L:
            tail = text[end:end+80]
            m = re.search(r'[.!?]\s', tail)
            if m:
                end += m.start() + 1
        chunk = text[start:end].strip()
        if len(chunk) > 20:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

# -------------------------
# Embeddings (OpenAI) with batching
# -------------------------
def get_openai_embeddings(texts: List[str], model:str) -> np.ndarray:
    # returns numpy float32 matrix (n, d)
    all_emb = []
    batch_size = 50
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        try:
            resp = openai.Embedding.create(model=model, input=batch)
            embs = [r["embedding"] for r in resp["data"]]
            all_emb.extend(embs)
        except Exception as e:
            st.error(f"Lỗi khi gọi Embedding API: {e}")
            raise e
    arr = np.array(all_emb, dtype="float32")
    return arr

# -------------------------
# Build index
# -------------------------
@st.cache_resource(show_spinner=False)
def build_vector_index(embeddings: np.ndarray, use_faiss_local=True):
    dim = embeddings.shape[1]
    meta = {"dim": dim}
    if use_faiss_local and FAISS_AVAILABLE:
        index = faiss.IndexFlatL2(dim)
        index.add(embeddings)
        meta["type"] = "faiss"
        return index, meta
    elif ANNOY_AVAILABLE:
        from annoy import AnnoyIndex
        idx = AnnoyIndex(dim, "angular")
        for i in range(embeddings.shape[0]):
            idx.add_item(i, embeddings[i].tolist())
        idx.build(10)
        meta["type"] = "annoy"
        return idx, meta
    else:
        # fallback: return None — we'll use brute force search via linear_kernel + stored embeddings+chunks
        meta["type"] = "bruteforce"
        return None, meta

# -------------------------
# Search
# -------------------------
def search_index(query_embedding: np.ndarray, index_obj, meta: Dict, top_k:int, embeddings_matrix: np.ndarray):
    if meta["type"] == "faiss":
        D, I = index_obj.search(query_embedding.astype("float32"), top_k)
        return I[0].tolist(), D[0].tolist()
    elif meta["type"] == "annoy":
        vec = query_embedding[0].tolist()
        ids, dists = index_obj.get_nns_by_vector(vec, top_k, include_distances=True)
        return ids, dists
    else:
        # brute force via cosine similarity
        sims = linear_kernel(query_embedding, embeddings_matrix).flatten()
        idxs = np.argsort(-sims)[:top_k]
        scores = sims[idxs].tolist()
        return idxs.tolist(), scores

# -------------------------
# UI: Upload files & process
# -------------------------
st.header("1) Upload tài liệu (PDF / DOCX / PPTX / XLSX / CSV / TXT)")
uploaded = st.file_uploader("Chọn file (multiple)", accept_multiple_files=True, type=["pdf","docx","pptx","xlsx","xls","csv","txt"])
process_btn = st.button("Xử lý & Index (tạo embeddings + index)")

if "docs" not in st.session_state:
    st.session_state["docs"] = []          # list of dicts: {file, chunk_id, text, source}
    st.session_state["chunks"] = []        # list of chunk text strings
    st.session_state["embeddings"] = None  # numpy array
    st.session_state["index"] = None
    st.session_state["index_meta"] = {}

if process_btn:
    if not uploaded:
        st.warning("Chưa có file nào được upload.")
    else:
        st.info("Đang trích xuất và chunk...")
        docs = []
        for f in uploaded:
            try:
                txt = extract_text_from_file(f)
            except Exception as e:
                st.error(f"Không đọc được file {f.name}: {e}")
                txt = ""
            chunks_local = chunk_text(txt, chunk_size=chunk_size, overlap=chunk_overlap)
            for i,ch in enumerate(chunks_local):
                docs.append({"source": f.name, "chunk_id": i, "text": ch})
        if not docs:
            st.warning("Không tạo được chunk từ file (file rỗng hoặc lỗi parse).")
        else:
            st.session_state["docs"] = docs
            st.session_state["chunks"] = [d["text"] for d in docs]
            st.success(f"Tạo được {len(st.session_state['chunks'])} chunk từ {len(uploaded)} file.")

            # embeddings
            if OPENAI_KEY:
                st.info("Đang gọi OpenAI để tạo embeddings (có thể mất vài chục giây)...")
                try:
                    emb = get_openai_embeddings(st.session_state["chunks"], model=DEFAULT_EMBED_MODEL)
                    st.session_state["embeddings"] = emb
                    st.success(f"Tạo embeddings xong — shape = {emb.shape}")
                    # build index
                    use_f = use_faiss and FAISS_AVAILABLE
                    index_obj, meta = build_vector_index(emb, use_faiss_local=use_f)
                    st.session_state["index"] = index_obj
                    st.session_state["index_meta"] = meta
                    st.success(f"Index đã tạo ({meta['type']}).")
                except Exception as e:
                    st.error(f"Embeddings lỗi: {e}")
                    st.session_state["embeddings"] = None
            else:
                st.info("Không có API key → không tạo embeddings. Sẽ dùng TF-IDF brute-force (nếu cần).")
                # build TF-IDF only for fallback search
                try:
                    tfidf = TfidfVectorizer(max_features=5000, stop_words='english')
                    X = tfidf.fit_transform(st.session_state["chunks"])
                    st.session_state["tfidf"] = tfidf
                    st.session_state["tfidf_matrix"] = X
                    st.session_state["index_meta"] = {"type":"tfidf"}
                    st.success("TF-IDF index đã sẵn sàng (fallback).")
                except Exception as e:
                    st.error(f"TF-IDF build lỗi: {e}")

# -------------------------
# Query & answer
# -------------------------
st.header("2) Truy vấn & Phân tích")
question = st.text_area("Nhập câu hỏi / yêu cầu phân tích", height=140)
ask_btn = st.button("Tìm và phân tích (tạo prompt gửi LLM)")

if ask_btn:
    if not question.strip():
        st.warning("Vui lòng nhập câu hỏi.")
    elif not st.session_state.get("chunks") or len(st.session_state["chunks"])==0:
        st.warning("Chưa index dữ liệu — nhấn 'Xử lý & Index' trước.")
    else:
        # embed query
        if OPENAI_KEY and st.session_state.get("embeddings") is not None:
            try:
                q_emb = get_openai_embeddings([question], model=DEFAULT_EMBED_MODEL)
            except Exception as e:
                st.error(f"Lỗi tạo embedding cho query: {e}")
                q_emb = None
        else:
            q_emb = None

        # search
        idxs = []
        scores = []
        meta = st.session_state.get("index_meta", {"type":"bruteforce"})
        if q_emb is not None and (meta["type"] in ["faiss","annoy","bruteforce"]):
            index_obj = st.session_state.get("index")
            emb_matrix = st.session_state.get("embeddings")
            try:
                ids, dists = search_index(q_emb, index_obj, meta, top_k, emb_matrix)
                idxs = ids
                scores = dists
            except Exception as e:
                st.error(f"Lỗi tìm kiếm trong index: {e}")
        else:
            # fallback to TF-IDF cosine similarity
            if "tfidf_matrix" in st.session_state:
                q_vec = st.session_state["tfidf"].transform([question])
                sims = linear_kernel(q_vec, st.session_state["tfidf_matrix"]).flatten()
                idxs = np.argsort(-sims)[:top_k].tolist()
                scores = sims[idxs].tolist()
            else:
                st.error("Không có index khả dụng để tìm kiếm (cần tạo embeddings hoặc TF-IDF).")

        # prepare context snippets
        snippets = []
        for rank, idx in enumerate(idxs):
            try:
                doc = st.session_state["docs"][idx]
                score = scores[rank] if rank < len(scores) else None
                header = f"[{rank+1}] {doc['source']} (chunk {doc['chunk_id']}) score={score:.4f}"
                snippets.append(header + "\n" + doc["text"])
            except Exception:
                continue

        # build prompt (careful about context length)
        max_chars_context = 6000
        ctx = ""
        for s in snippets:
            if len(ctx) + len(s) > max_chars_context:
                break
            ctx += s + "\n\n"

        prompt = (
            "Bạn là kỹ sư bảo dưỡng máy bay. Dưới đây là các đoạn trích từ hồ sơ/báo cáo (relevant context):\n\n"
            f"{ctx}\n\n"
            "Câu hỏi của người dùng:\n"
            f"{question}\n\n"
            "Yêu cầu: Hãy trả lời bằng tiếng Việt. Bao gồm: (1) Tóm tắt ngắn các phát hiện quan trọng; "
            "(2) Nguyên nhân khả dĩ; (3) Hành động/kiểm tra đề xuất; (4) Mức độ khẩn cấp (high/medium/low). "
            "Trình bày có tiêu đề và bullet points."
        )

        st.subheader("Prompt (preview)")
        st.code(prompt[:4000] + ("...\n(đã cắt)" if len(prompt) > 4000 else ""), language="text")

        # call LLM for analysis
        if OPENAI_KEY:
            try:
                with st.spinner("Gọi LLM để phân tích (OpenAI)..."):
                    resp = openai.ChatCompletion.create(
                        model=DEFAULT_CHAT_MODEL,
                        messages=[
                            {"role":"system","content":"You are a helpful aviation maintenance analyst. Keep answers concise and technical in Vietnamese."},
                            {"role":"user","content":prompt}
                        ],
                        max_tokens=800,
                        temperature=0.0,
                    )
                    answer = resp["choices"][0]["message"]["content"].strip()
                    st.subheader("Kết quả (AI)")
                    st.markdown(answer)
            except Exception as e:
                st.error(f"Lỗi gọi ChatCompletion: {e}")
        else:
            st.info("Không có API key → trả về kết quả phân tích đơn giản (fallback).")
            summary_lines = [f"- {s[:300].replace('\\n',' ')}" for s in snippets]
            fallback = "Tóm tắt (fallback):\n" + "\n".join(summary_lines)
            st.text(fallback)

        # allow download
        result_txt = f"QUESTION:\n{question}\n\nCONTEXT:\n{ctx}\n\nPROMPT:\n{prompt}\n\nANSWER:\n{answer if OPENAI_KEY else fallback}"
        st.download_button("Tải kết quả (TXT)", result_txt, file_name="analysis_result.txt", mime="text/plain")

# -------------------------
# Utilities: save/load index (optional)
# -------------------------
st.header("3) (Tùy chọn) Lưu / Load index (local file trong environment)")
col1, col2 = st.columns(2)
with col1:
    if st.button("Lưu index & metadata (file index.faiss + meta.pkl)"):
        if st.session_state.get("index") is None and st.session_state.get("embeddings") is None:
            st.warning("Chưa có index để lưu.")
        else:
            try:
                # save embeddings+meta+docs
                with open("meta_docs.pkl", "wb") as f:
                    pickle.dump({"docs": st.session_state["docs"], "meta": st.session_state.get("index_meta")}, f)
                if st.session_state["index_meta"].get("type") == "faiss" and FAISS_AVAILABLE:
                    faiss.write_index(st.session_state["index"], "index.faiss")
                elif st.session_state["index_meta"].get("type") == "annoy" and ANNOY_AVAILABLE:
                    st.session_state["index"].save("index.ann")
                # save embeddings for brute force
                if st.session_state.get("embeddings") is not None:
                    np.save("embeddings.npy", st.session_state["embeddings"])
                st.success("Đã lưu các file index (trong filesystem của app).")
            except Exception as e:
                st.error(f"Lỗi lưu file: {e}")
with col2:
    uploaded_index = st.file_uploader("Upload file index zip (optional)", type=["zip","pkl","npy","faiss","ann"])
    if st.button("Load index từ file upload"):
        st.warning("Tính năng load index upload chưa triển khai chi tiết — bạn có thể triển khai theo định dạng đã lưu.")

# Footer notes
st.markdown("---")
st.markdown(
    "## Ghi chú vận hành\n"
    "- File system trên nhiều cloud (Streamlit Cloud) là ephemeral giữa deploys — để lưu index lâu dài cần external storage (S3/Pinecone/Weaviate).\n"
    "- Nếu `faiss-cpu` không cài được trên môi trường, code sẽ fallback sang Annoy hoặc TF-IDF.\n"
    "- Hạn chế số chunk đưa vào prompt để tiết kiệm token & tiền. Điều chỉnh `top_k` và `chunk_size`.\n"
    "- Với tài liệu quét (scan image PDF), cần OCR (tesseract/pytesseract) trước khi extract text.\n\n"
    "Muốn tôi lọc thêm (ví dụ chỉ index cột 'description' trong Excel, hoặc thêm OCR), nói tôi biết nhé!"
)
